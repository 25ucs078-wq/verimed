import os
import sys

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not installed. Attempting installation...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def build_presentation():
    prs = Presentation()
    
    # Slide Dimensions (16:9 widescreen layout standard)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Custom Theme Color Scheme (Dark space-navy theme matching VeriMed branding)
    COLOR_BG = RGBColor(7, 10, 24)       # Deep Space Navy #070A18
    COLOR_TEAL = RGBColor(45, 224, 194)  # Electric Teal #2DE0C2
    COLOR_RED = RGBColor(255, 85, 119)   # Coral Red #FF5577
    COLOR_WHITE = RGBColor(248, 250, 252)# Slate 50 #F8FAFC
    COLOR_MUTED = RGBColor(148, 163, 184)# Slate 400 #94A3B8
    COLOR_CARD = RGBColor(16, 19, 43)    # Card Navy #10132B
    
    # Helper to set solid slide background
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG
        
    # Helper to add text frame box
    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Space Grotesk"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL
        return title_box

    # ------------------ SLIDE 1: Title Slide ------------------
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    
    # Title & Subtitle box
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "VeriMed"
    p1.font.name = "Space Grotesk"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    
    # Colored tag text
    run = p1.add_run()
    run.text = " Intelligence"
    run.font.color.rgb = COLOR_TEAL
    
    p2 = tf.add_paragraph()
    p2.text = "Scan. Verify. Stay Safe. — AI-Powered Counterfeit Medicine Countermeasures"
    p2.font.name = "Inter"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_MUTED
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "Product Proposal & Ecosystem Showcase"
    p3.font.name = "Inter"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_TEAL
    p3.space_before = Pt(50)

    # ------------------ SLIDE 2: The Core Challenge ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, "The Counterfeit Drug Epidemic")
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Problem description bullet items
    points = [
        ("Global Threat Vector", "Up to 10% of medical products in low- and middle-income countries are falsified or substandard (WHO estimates)."),
        ("Severe Healthcare Impact", "Responsible for hundreds of thousands of preventable deaths annually, including pediatric pneumonia and malaria cases."),
        ("Connectivity Gaps", "Traditional verification networks fail in remote regions (Sub-Saharan Africa, rural South Asia) due to offline environments."),
        ("Advanced Imitations", "Holograms and packaging are increasingly easy for counterfeiters to visually clone, requiring algorithmic vision verification.")
    ]
    
    for idx, (title, desc) in enumerate(points):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = f"•  {title}: "
        p.font.name = "Inter"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_MUTED
        
    # Screenshot container placeholder right
    rect = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_CARD
    rect.line.color.rgb = COLOR_TEAL
    rect.line.width = Pt(1.5)
    
    tf_rect = rect.text_frame
    tf_rect.word_wrap = True
    p_rect = tf_rect.paragraphs[0]
    p_rect.text = "[ DRAG & DROP SCREENSHOT: REGIONAL NEWS / ALERTS HUB ]"
    p_rect.font.name = "Space Grotesk"
    p_rect.font.size = Pt(13)
    p_rect.font.color.rgb = COLOR_MUTED
    p_rect.alignment = PP_ALIGN.CENTER

    # ------------------ SLIDE 3: The VeriMed Solution ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, "The VeriMed Solution Architecture")
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    features = [
        ("AI Packaging Print Vision", "Evaluates logo displacements, typographical kerning anomalies, and color spectrum deviations in real-time.", COLOR_TEAL),
        ("SmartID Cryptographic QR Decrypts", "Decodes DCGI-standard QR parameters and verifies digital signatures utilizing cached manufacturer public keys.", COLOR_TEAL),
        ("Offline-First Synchronizer Center", "Saves verification outboxes locally to SQLite / LocalStorage, and automatically syncs when network connectivity returns.", COLOR_TEAL),
        ("Dynamic Threat Maps & News Hub", "Gathers scanning metadata to dynamically map local threat indices and alerts operators in high-risk zones.", COLOR_TEAL)
    ]
    
    for idx, (title, desc, color) in enumerate(features):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = f"0{idx+1}. {title}\n"
        p.font.name = "Space Grotesk"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)
        if idx > 0:
            p.space_before = Pt(16)
            
        run = p.add_run()
        run.text = desc
        run.font.name = "Inter"
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = COLOR_MUTED

    # ------------------ SLIDE 4: Product Showcase: Dashboard ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, "Product Showcase: Integrity Dashboard")
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Central Command Hub"
    p.font.name = "Space Grotesk"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(14)
    
    points = [
        "Dynamic counts displaying total scans conducted and counterfeits blocked.",
        "Regional Threat Index widgets showing active threat levels and trend metrics.",
        "Recent activity logs providing immediate audits on verified and suspect batches.",
        "Responsive sidebar menu converting into a floating glass tab bar on mobile layouts."
    ]
    
    for pt in points:
        p2 = tf.add_paragraph()
        p2.text = f"•  {pt}"
        p2.font.name = "Inter"
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_after = Pt(10)
        
    # Screenshot placeholder
    rect = slide.shapes.add_shape(
        1, Inches(5.8), Inches(1.8), Inches(6.8), Inches(4.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_CARD
    rect.line.color.rgb = COLOR_TEAL
    rect.line.width = Pt(1.5)
    
    tf_rect = rect.text_frame
    tf_rect.word_wrap = True
    p_rect = tf_rect.paragraphs[0]
    p_rect.text = "[ DRAG & DROP SCREENSHOT: INTEGRITY DASHBOARD VIEW ]"
    p_rect.font.name = "Space Grotesk"
    p_rect.font.size = Pt(13)
    p_rect.font.color.rgb = COLOR_MUTED
    p_rect.alignment = PP_ALIGN.CENTER

    # ------------------ SLIDE 5: Product Showcase: Scanner ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, "Product Showcase: Verification Viewfinder")
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Real-time Verification Scanner"
    p.font.name = "Space Grotesk"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(14)
    
    points = [
        "Switches between print packaging vision analysis and DCGI standard barcode scan modes.",
        "Sweeping laser animations guide the operator to align components inside the viewfinder frame.",
        "Developer simulator triggers facilitate comprehensive testing without native hardware cameras.",
        "Instantly triggers detailed reports or alerts if print or signature tolerances fail."
    ]
    
    for pt in points:
        p2 = tf.add_paragraph()
        p2.text = f"•  {pt}"
        p2.font.name = "Inter"
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_after = Pt(10)
        
    # Screenshot placeholder
    rect = slide.shapes.add_shape(
        1, Inches(5.8), Inches(1.8), Inches(6.8), Inches(4.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_CARD
    rect.line.color.rgb = COLOR_TEAL
    rect.line.width = Pt(1.5)
    
    tf_rect = rect.text_frame
    tf_rect.word_wrap = True
    p_rect = tf_rect.paragraphs[0]
    p_rect.text = "[ DRAG & DROP SCREENSHOT: SCANNER VIEWFINDER VIEW ]"
    p_rect.font.name = "Space Grotesk"
    p_rect.font.size = Pt(13)
    p_rect.font.color.rgb = COLOR_MUTED
    p_rect.alignment = PP_ALIGN.CENTER

    # ------------------ SLIDE 6: Offline-First Synchronization ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, "Offline-First Sync Strategy")
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    points = [
        ("Local Cache Mirroring", "Uses LocalStorage (web) / SQLite (native mobile) databases to cache up to 100 batch files and manufacturer keys within a 50km radius of the operator."),
        ("Queued Outbox Pipeline", "If disconnected, scans and reports are stored locally. Sync Center shows outbox status and triggers dynamic uploads once back online."),
        ("Dynamic Key Fetching", "A 7-day expiration index downloads updated cryptographic signature prefixes to verify batches offline without server lookups.")
    ]
    
    for idx, (title, desc) in enumerate(points):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = f"•  {title}: "
        p.font.name = "Inter"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_MUTED
        
    # Screenshot placeholder right
    rect = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_CARD
    rect.line.color.rgb = COLOR_TEAL
    rect.line.width = Pt(1.5)
    
    tf_rect = rect.text_frame
    tf_rect.word_wrap = True
    p_rect = tf_rect.paragraphs[0]
    p_rect.text = "[ DRAG & DROP SCREENSHOT: SYNC & CACHE CENTER VIEW ]"
    p_rect.font.name = "Space Grotesk"
    p_rect.font.size = Pt(13)
    p_rect.font.color.rgb = COLOR_MUTED
    p_rect.alignment = PP_ALIGN.CENTER

    # ------------------ SLIDE 7: Technical Stack Summary ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    add_title(slide, "Technical Stack & Scalability")
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    stack = [
        ("High-Performance Backend", "FastAPI, SQLAlchemy 2.0 ORM, AsyncSession, Uvicorn server, Pydantic V2 validations, PostgreSQL, and Redis cache databases."),
        ("Modular Service Layer", "Decoupled services for EfficientNet-B0 CV simulations, QR cryptographical decodes, and geopolitical risk engines."),
        ("Responsive Web Dashboard", "Vite + React + Tailwind CSS client with Lucide icons. Desktop navigation converts into a bottom tab bar on mobile layouts."),
        ("Security Standards", "Auto-documented REST APIs under Swagger (/docs) with full HTTP Exception handling and token-based signature verification.")
    ]
    
    for idx, (title, desc) in enumerate(stack):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = f"•  {title}: "
        p.font.name = "Inter"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(6)
        if idx > 0:
            p.space_before = Pt(10)
            
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_MUTED

    # ------------------ SLIDE 8: Summary / Q&A ------------------
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "VeriMed: Scan. Verify. Stay Safe."
    p1.font.name = "Space Grotesk"
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEAL
    
    p2 = tf.add_paragraph()
    p2.text = "Establishing a trust boundary for clinical supply chains in remote regions."
    p2.font.name = "Inter"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "Thank you! Open for Questions."
    p3.font.name = "Inter"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(60)

    # Save presentation
    filepath = "VeriMed_Product_Pitch.pptx"
    prs.save(filepath)
    print(f"Presentation saved successfully to: {os.path.abspath(filepath)}")

if __name__ == "__main__":
    build_presentation()
